import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import streamlit as st
import matplotlib.font_manager
from streamlit_option_menu import option_menu

# charts
from charts_bar import add_bar_slide

# def load_ppt_template():
#     selected_template = st.selectbox('TEMPLATES', options=['Template-1', 'Template-2'])

#     if selected_template == 'Template-1':
#         return Presentation('Templates/Template-1.pptx')
    
#     if selected_template == 'Template-2':
#         return Presentation('Templates/Template-2.pptx')


# def create_bar_chart(df, slide, idx):

#     # Add Category
#     chart_data = CategoryChartData()
#     chart_data.categories = df['Date']

#     # Add Series
#     df_headers = df.columns.tolist()[1:]
#     for df_header in df_headers:
#         chart_data.add_series(df_header, df[df_header])
    
#     # Insert chart to the placeholder
#     chart_placeholder = slide.placeholders[idx]
#     chart_frame = chart_placeholder.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data)
#     _chart = chart_frame.chart

#     return _chart


# def get_available_fonttypeface():

#     # Get a list of all available font names
#     available_fonts = sorted(set(f.name for f in matplotlib.font_manager.fontManager.ttflist))
    
#     return available_fonts


def add_pie_slide(df, prs, grouping):

    if grouping == 'Standard':
        grouptype = XL_CHART_TYPE.PIE
        layoutname = 'PieChartStandard'
    if grouping == 'Doughnut':
        grouptype = XL_CHART_TYPE.DOUGHNUT
        layoutname = 'PieChartDoughnut'
    if grouping == 'Exploded Standard':
        grouptype = XL_CHART_TYPE.PIE_EXPLODED
        layoutname = 'PieChartxStandard'
    if grouping == 'Exploded Doughnut':
        grouptype = XL_CHART_TYPE.DOUGHNUT_EXPLODED
        layoutname = 'PieChartxDoughnut'

    # get the index number of the Chart Placeholder from the slide named BarChart
    for layout in prs.slide_layouts:
        if layout.name == layoutname:
            slide = prs.slides.add_slide(layout)
            for placeholder in layout.placeholders:
                if 'Chart Placeholder' in placeholder.name:
                    idx = placeholder.placeholder_format.idx
    
    chart_data = add_chart_data()
      
    # Add pie chart to the slide
    chart_placeholder = slide.placeholders[idx]
    chart_frame = chart_placeholder.insert_chart(grouptype, chart_data)
    _chart = chart_frame.chart

    chart_details(_chart, False, False)


def add_line_slide(df, prs):

    # get the index number of the Chart Placeholder from the slide named BarChart
    for layout in prs.slide_layouts:
        if layout.name == 'LineChart':
            slide = prs.slides.add_slide(layout)
            for placeholder in layout.placeholders:
                if 'Chart Placeholder' in placeholder.name:
                    idx = placeholder.placeholder_format.idx

    chart_data = add_chart_data()
      
    # Add pie chart to the slide
    chart_placeholder = slide.placeholders[idx]
    chart_frame = chart_placeholder.insert_chart(XL_CHART_TYPE.LINE, chart_data)
    _chart = chart_frame.chart

    chart_details(_chart, True, True)


def add_chart_data():

    # Define chart data
    chart_data = CategoryChartData()
    chart_data.categories = df[df.columns.tolist()[0]]
    
    df_headers = df.columns.tolist()[1:]
    for df_header in df_headers:
        chart_data.add_series(df_header, df[df_header])
    
    return chart_data


def chart_details(_chart, cat_axis, value_axis):

    # chart title
    _chart.has_title = True
    _chart.chart_title.text_frame.text = "Place Chart Title Here"
    try:
        _chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        _chart.chart_title.text_frame.paragraphs[0].font.name = 'Arial'
    except:
        pass
    
    if cat_axis == True:
        # category axis
        _chart.category_axis.axis_title.text_frame.text = df.columns.tolist()[0]
        try:
            _chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(18)
            _chart.category_axis.axis_title.text_frame.paragraphs[0].font.name = 'Arial'
        except:
            pass
    
    if value_axis == True:
        # value axis
        _chart.value_axis.axis_title.text_frame.text = "Count"
        try:
            _chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(18)
            _chart.value_axis.axis_title.text_frame.paragraphs[0].font.name = 'Arial'
        except:
            pass
    
    # legend
    _chart.has_legend = True
    _chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    _chart.legend.include_in_layout = False
    _chart.legend.font.size = Pt(8)
    _chart.legend.font.bold = True
    _chart.legend.font.name = 'Arial'

    # data labels
    for series in _chart.series:
        series.has_data_labels = True  # Enable data labels
        series.data_labels.show_value = True  # Show values on the labels
        # series.data_labels.show_category_name = True
        series.data_labels.font.size = Pt(8)
        series.data_labels.font.bold = True
        series.data_labels.font.name = 'Arial'
        # series.data_labels.show_series_name = False
        # series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    
    # data table
    _chart.has_data_table = True
    # _chart.data_table.has_border_horizontal = True  # Add horizontal borders
    # _chart.data_table.has_border_vertical = True  # Add vertical borders
    # _chart.data_table.has_border_outline = True  # Add outline border








if __name__ == '__main__':
    
    with st.sidebar:
        chart_type = option_menu(
            menu_title='Chart Type',
            options=['Bar', 'Pie', 'Line']
        )

        if chart_type == 'Bar':
            bar_orientation = option_menu(
                menu_title='Bar Orientation',
                options=['Horizontal', 'Vertical']
            )
            sub_type = option_menu(
                menu_title='Grouping',
                options=['Stacked', 'Clustered', '100 Percent Stacked']
            )
        
        if chart_type == 'Pie':
            sub_type = option_menu(
                menu_title='Grouping',
                options=['Standard', 'Doughnut', 'Exploded Standard', 'Exploded Doughnut']
                # options=['Standard', 'Doughnut']
            )
    data_csv = st.file_uploader('Data', type='csv', label_visibility='hidden')

    if data_csv not in [None, '']:

        # create presentation from the selected template
        prs = Presentation('Templates/Template-3.pptx')

        # load data
        df = pd.read_csv(data_csv)
        
        if chart_type == 'Bar':
            if sub_type == 'Clustered':
                add_bar_slide(df, prs, 'Clustered', bar_orientation)               
            
            if sub_type == 'Stacked':
                add_bar_slide(df, prs, 'Stacked', bar_orientation)
            
            if sub_type == '100 Percent Stacked':
                add_bar_slide(df, prs, '100 Percent Stacked', bar_orientation)
        else:
            pass
        
        if chart_type == 'Pie':
            if sub_type == 'Standard':
                add_pie_slide(df, prs, 'Standard')               
            
            elif sub_type == 'Doughnut':
                add_pie_slide(df, prs, 'Doughnut')

            elif sub_type == 'Exploded Standard':
                add_pie_slide(df, prs, 'Exploded Standard')

            elif sub_type == 'Exploded Doughnut':
                add_pie_slide(df, prs, 'Exploded Doughnut')
        else:
            pass
        
        if chart_type == 'Line':
            add_line_slide(df, prs)

                
        # save the presentation
        btn_save = st.button('Save')
        if btn_save:
            prs.save('Output/output_presentation4.pptx')

            result_file = open('Output/output_presentation4.pptx', 'rb')
            st.success(f':red[NOTE:] Downloaded file will go to the :red[Downloads Folder]')
            st.download_button(label='📥 Download Cleaned Raw', data=result_file ,file_name= f'ppt_test_charts.pptx')


    


    # csv_file = st.file_uploader('Upload File', type=['csv'], key='file_csv')
    
    # if st.session_state['file_csv'] not in [None, '']:
    #     df = pd.read_csv(st.session_state['file_csv'])
    #     with st.expander('VIEW DATAFRAME'):
    #         st.dataframe(df, use_container_width=True, hide_index=True)
    
    #     # load template
    #     _layouts = {}
    #     _layout_names = []

    #     col11, col12 = st.columns(2)
        
    #     with col11:
    #         prs = load_ppt_template()

    #     with col12:
    #         for layout in prs.slide_layouts:
    #             _layouts[layout.name] = layout
    #             _layout_names.append(layout.name)
        
    #         selected_layout = st.selectbox('LAYOUTS', options=_layout_names)
    #         my_layout = _layouts[selected_layout]
    #         for placeholder in my_layout.placeholders:
    #             if 'Chart Placeholder' in placeholder.name:
    #                 # st.write(f'{placeholder.name} - {placeholder.placeholder_format.idx}')
    #                 idx = placeholder.placeholder_format.idx
        
    #     slide = prs.slides.add_slide(_layouts[selected_layout])

    #     chart = create_bar_chart(df, slide, idx)

    #     with st.expander('CHART DETAILS'):
    #         col21, col22, col23 = st.columns(3, border=True)

    #         # chart title
    #         with col21:
    #             cht_title = st.text_input('Chart Title')
    #             cht_title_font_size = st.number_input('Font Size', min_value=1, max_value=100, step=1, value=8, help='Size of the chart title font, default is 8')
    #             cht_title_font_name = st.selectbox('Font Name', options=get_available_fonttypeface())
                
    #             chart.has_title = True
    #             if cht_title in [None, '']:
    #                 chart.chart_title.text_frame.text = 'BAR CHART TITLE'
    #             else:
    #                 chart.chart_title.text_frame.text = cht_title

    #             try:
    #                 chart.chart_title.text_frame.paragraphs[0].font.size = Pt(cht_title_font_size)
    #                 chart.chart_title.text_frame.paragraphs[0].font.name = cht_title_font_name
    #             except:
    #                 pass

    #     btn_save = st.button('CREATE CHART')
    #     if btn_save:
    #         prs.save('Output/output_presentation3.pptx')


    # _title = slide.shapes.title
    # _title.text = 'My Sample Chart'

    # # Add Category
    # chart_data = CategoryChartData()
    # chart_data.categories = df['Date']

    # # Add Series
    # df_headers = df.columns.tolist()[1:]
    # for df_header in df_headers:
    #     chart_data.add_series(df_header, df[df_header])    
    
    # # Insert chart to the placeholder
    # chart_placeholder = slide.placeholders[_placeholderindex]
    # chart_frame = chart_placeholder.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data)
    # chart = chart_frame.chart

    # # Customize chart title
    # chart.has_title = True
    # chart.chart_title.text_frame.text = "Motorcycle Sales Over Time"
    # try:
    #     chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
    #     chart.chart_title.text_frame.paragraphs[0].font.name = 'Arial'
    # except:
    #     pass

    # # Customize category axis
    # chart.category_axis.axis_title.text_frame.text = "Date"
    # try:
    #     chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(8)
    #     chart.category_axis.axis_title.text_frame.paragraphs[0].font.name = 'Arial'
    # except:
    #     pass

    # # Customize category axis tick labels
    #     tick_labels_xml = chart.category_axis._element.find(".//c:txPr", namespaces={"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"})
        
    #     if tick_labels_xml is not None:
    #         print('Yes')


    #     chart.category_axis.tick_labels.font.bold = True
    #     chart.category_axis.tick_labels.font.color.rgb = RGBColor(100, 0, 0)
    #     chart.value_axis.tick_labels.rotation = 90

    # # Customize series axis
    # chart.value_axis.axis_title.text_frame.text = "Sales Count"
    # try:
    #     chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(8)
    #     chart.value_axis.axis_title.text_frame.paragraphs[0].font.name = 'Arial'
    # except:
    #     pass

    # # Customize gap between bars
    # chart.plots[0].gap_width = 50

    # # Apply custom colors to series
    # series_colors = [RGBColor(255, 0, 0), RGBColor(0, 128, 0), RGBColor(0, 0, 255), RGBColor(255, 165, 0), RGBColor(128, 0, 128)]

    # # chart transparency
    # # chart_frame.fill.solid()
    # # chart_frame.fill.fore_color.transparency = 0.0
    # # chart.plot_area.format.fill.solid()
    # # chart.plot_area.format.fill.fore_color.transparency = 0.0

    # # Customize legend
    # chart.has_legend = True
    # chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    # chart.legend.include_in_layout = False
    # chart.legend.font.size = Pt(8)
    # chart.legend.font.bold = True
    # chart.legend.font.name = 'Arial'
    
    

    # for layout in prs.slide_layouts:
    #     try:
    #         prs.slide_layouts.remove(layout)
    #     except:
    #         pass
    
    # prs.save('Output/output_presentation3.pptx')