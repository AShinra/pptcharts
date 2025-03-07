import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import streamlit as st
import matplotlib.font_manager
from streamlit_option_menu import option_menu
import json


def get_available_fonttypeface():
    with open('fonts.json', 'r') as file:
        fonts = json.load(file)                                                                                                     
    return fonts



def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def add_bar_slide(df, prs, grouping, bar_orientation):

    if grouping == 'Clustered':
        if bar_orientation == 'Vertical':
            grouptype = XL_CHART_TYPE.COLUMN_CLUSTERED
        elif bar_orientation == 'Horizontal':
            grouptype = XL_CHART_TYPE.BAR_CLUSTERED
        layoutname = 'BarChartClustered'
    if grouping == 'Stacked':
        if bar_orientation == 'Vertical':
            grouptype = XL_CHART_TYPE.COLUMN_STACKED
        elif bar_orientation == 'Horizontal':
            grouptype = XL_CHART_TYPE.BAR_STACKED
        layoutname = 'BarChartStacked'
    if grouping == '100 Percent Stacked':
        if bar_orientation == 'Vertical':
            grouptype = XL_CHART_TYPE.COLUMN_STACKED_100
        elif bar_orientation == 'Horizontal':
            grouptype = XL_CHART_TYPE.BAR_STACKED_100
        layoutname = 'BarChart100Stacked'

    # get the index number of the Chart Placeholder from the slide named BarChart
    for layout in prs.slide_layouts:
        if layout.name == layoutname:
            slide = prs.slides.add_slide(layout)
            for placeholder in layout.placeholders:
                if 'Chart Placeholder' in placeholder.name:
                    idx = placeholder.placeholder_format.idx

    # Define chart data
    chart_data = CategoryChartData()
    chart_data.categories = df[df.columns.tolist()[0]]
    
    df_headers = df.columns.tolist()[1:]
    for df_header in df_headers:
        chart_data.add_series(df_header, df[df_header])
    
    # Insert chart to the placeholder
    chart_placeholder = slide.placeholders[idx]
    chart_frame = chart_placeholder.insert_chart(grouptype, chart_data)
    _chart = chart_frame.chart



    # -------------- chart details -------------------
    tab1, tab2, tab3, tab4, tab5 = st.tabs(['VIEW DATA', 'CHART TITLE', 'CATEGORY AXIS', 'VALUE AXIS', 'LEGEND'])

    # view dataframe
    with tab1:
        st.dataframe(df, hide_index='hidden', use_container_width=True)
        
    # chart title
    available_fonts = get_available_fonttypeface()

    cht_title_dict = {}
    with tab2:
        cht_title = st.text_input('Chart Title', key='cht_title', placeholder='Input Chart Title Here', label_visibility='collapsed')
        if cht_title in [None, '']:
            cht_title_dict['text'] = 'No Chart Title'
        else:
            cht_title_dict['text'] = cht_title

        col_title11, col_title12, col_title13 = st.columns([0.75, 0.15, 0.10])
        with col_title11:
            cht_title_dict['font_name'] = st.selectbox('Font Name', options=available_fonts, key='cht_title_font_name')
        with col_title12:
            cht_title_dict['font_size'] = st.number_input('Font Size', min_value=1, max_value=100, value=10, step=1, key='cht_title_font_size')
        with col_title13:
            cht_title_font_color = st.color_picker('Color', key='cht_title_font_color')
            rgb_color = hex_to_rgb(cht_title_font_color)
            cht_title_dict['font_color'] = rgb_color

        col_title21, col_title22, col_title23, col_title24 = st.columns([0.15, 0.15, 0.15, 0.55])
        with col_title21:
            if st.checkbox('Bold', key='title_bold'):
                cht_title_dict['bold'] = True
            else:
                cht_title_dict['bold'] = False
        with col_title22:
            if st.checkbox('Italic', key='title_italic'):
                cht_title_dict['italic'] = True
            else:
                cht_title_dict['italic'] = False
        with col_title23:
            if st.checkbox('Underline', key='title_underline'):
                cht_title_dict['underline'] = True
            else:
                cht_title_dict['underline'] = False

    # category axis
    cht_category_axis_dict = {}
    with tab3:
        cat_label = st.text_input('Category Label', key='cat_label', placeholder='Input Category Label Here', label_visibility='collapsed')
        if cat_label in [None, '']:
            cht_category_axis_dict['text'] = df.columns.tolist()[0]
        else:   
            cht_category_axis_dict['text'] = cat_label
        
        col_cat11, col_cat12, col_cat13 = st.columns([0.75, 0.15, 0.10])
        with col_cat11:
            cht_category_axis_dict['font_name'] = st.selectbox('Font Name', options=available_fonts, key='cht_category_font_name')
        with col_cat12:
            cht_category_axis_dict['font_size'] = st.number_input('Font Size', min_value=1, max_value=100, value=10, step=1, key='cht_category_font_size')
        with col_cat13:
            cht_cat_axis_ftcolor = st.color_picker('Color', key='cht_category_font_color')
            rgb_color = hex_to_rgb(cht_cat_axis_ftcolor)
            cht_category_axis_dict['font_color'] = rgb_color

        
        col_cat21, col_cat22, col_cat23, col_cat24 = st.columns([0.15, 0.15, 0.15, 0.55])
        with col_cat21:
            if st.checkbox('Bold', key='cat_bold'):
                cht_category_axis_dict['bold'] = True
            else:
                cht_category_axis_dict['bold'] = False
        with col_cat22:
            if st.checkbox('Italic', key='cat_italic'):
                cht_category_axis_dict['italic'] = True
            else:
                cht_category_axis_dict['italic'] = False
        with col_cat23:
            if st.checkbox('Underline', key='cat_underline'):
                cht_category_axis_dict['underline'] = True
            else:
                cht_category_axis_dict['underline'] = False
        
    
    # value axis
    cht_value_axis_dict = {}
    with tab4:
        if grouping == '100 Percent Stacked':
            cht_value_axis_dict['text'] = 'Percent'
        else:
            cht_value_axis_dict['text'] = st.text_input('Value Label', key='val_label', placeholder='Input Value Label Here')
        
        col_val1, col_val2, col_val3 = st.columns([0.75, 0.15, 0.10])
        with col_val1:
            cht_value_axis_dict['font_name'] = st.selectbox('Font Name', options=available_fonts, key='cht_value_font_name')
        with col_val2:
            cht_value_axis_dict['font_size'] = st.number_input('Font Size', min_value=1, max_value=100, value=10, step=1, key='cht_value_font_size')
        with col_val3:
            cht_value_axis_ftcolor = st.color_picker('Font Color', key='cht_value_font_color')
            rgb_color = hex_to_rgb(cht_value_axis_ftcolor)
            cht_value_axis_dict['font_color'] = rgb_color
        
    
    # legend
    cht_legend_dict = {}
    with tab5:
        col_legend11, col_legend12, col_legend13, col_legend14 = st.columns(4)
        with col_legend11:
            cht_legend_dict['font_name'] = st.selectbox('Font Name', options=available_fonts, key='cht_legend_font_name')
        with col_legend12:
            cht_legend_dict['font_size'] = st.number_input('Font Size', min_value=1, max_value=100, value=10, step=1, key='cht_legend_font_size')
        with col_legend13:
            if st.checkbox('Bold', key='cht_legend_bold'):
                cht_legend_dict['bold'] = True
            else:
                cht_legend_dict['bold'] = False
        with col_legend14:
            if st.checkbox('Italic', key='cht_legend_italic'):
                cht_legend_dict['italic'] = True
            else:
                cht_legend_dict['italic'] = False
        col_legend21, col_legend22 = st.columns(2) 
        with col_legend21:
            legend_loc = st.selectbox('Location', options=['Bottom', 'Top', 'Corner', 'Left', 'Right'])
            if legend_loc == 'Bottom':
                cht_legend_dict['location'] = XL_LEGEND_POSITION.BOTTOM
            elif legend_loc == 'Top':
                cht_legend_dict['location'] = XL_LEGEND_POSITION.TOP
            elif legend_loc == 'Corner':
                cht_legend_dict['location'] = XL_LEGEND_POSITION.CORNER
            elif legend_loc == 'Left':
                cht_legend_dict['location'] = XL_LEGEND_POSITION.LEFT
            elif legend_loc == 'Right':
                cht_legend_dict['location'] = XL_LEGEND_POSITION.RIGHT
        with col_legend22:
            legend_bool = st.checkbox('Show')
            if legend_bool:
                cht_legend_dict['visible'] = True
            else:
                cht_legend_dict['visible'] = False



    chart_details(df, _chart, cht_title_dict, cht_category_axis_dict, cht_value_axis_dict, cht_legend_dict)

    return


def chart_details(df, _chart, cht_title_dict, cht_category_axis_dict, cht_value_axis_dict, cht_legend_dict):

    # chart title
    _chart.has_title = True
    cht_title = _chart.chart_title.text_frame
    cht_title.text = cht_title_dict['text']
    try:
        cht_title.paragraphs[0].font.size = Pt(cht_title_dict['font_size'])
        cht_title.paragraphs[0].font.name = cht_title_dict['font_name']
        cht_title.paragraphs[0].font.color.rgb = RGBColor(*cht_title_dict['font_color'])
        cht_title.paragraphs[0].font.bold = cht_title_dict['bold']
        cht_title.paragraphs[0].font.italic = cht_title_dict['italic']
        cht_title.paragraphs[0].font.underline = cht_title_dict['underline']
    except:
        pass
    
    # category axis
    chart_category_axis = _chart.category_axis.axis_title.text_frame
    chart_category_axis.text = df.columns.tolist()[0]
    try:
        chart_category_axis.paragraphs[0].font.size = Pt(cht_category_axis_dict['font_size'])
        chart_category_axis.paragraphs[0].font.name = cht_category_axis_dict['font_name']
        chart_category_axis.paragraphs[0].font.color.rgb = RGBColor(*cht_category_axis_dict['font_color'])
    except:
        pass
    
    # value axis
    chart_value_axis = _chart.value_axis.axis_title.text_frame
    if cht_value_axis_dict['text'] in [None, '']:
        chart_value_axis.text = 'Count'
    else:    
        chart_value_axis.text = cht_value_axis_dict['text']
    
    try:
        chart_value_axis.paragraphs[0].font.size = Pt(cht_value_axis_dict['font_size'])
        chart_value_axis.paragraphs[0].font.name = cht_value_axis_dict['font_name']
        chart_value_axis.paragraphs[0].font.color.rgb = RGBColor(*cht_value_axis_dict['font_color'])
    except:
        pass

    # legend
    if cht_legend_dict['visible'] == True:
        _chart.has_legend = True
        _chart.legend.include_in_layout = False
        _chart.legend.position = cht_legend_dict['location']
        _chart.legend.font.name = cht_legend_dict['font_name']
        _chart.legend.font.size = Pt(cht_legend_dict['font_size'])
        _chart.legend.font.bold = cht_legend_dict['bold']
        _chart.legend.font.italic = cht_legend_dict['italic']
    else:
        pass
    

    # # data labels
    # for series in _chart.series:
    #     series.has_data_labels = True  # Enable data labels
    #     series.data_labels.show_value = True  # Show values on the labels
    #     # series.data_labels.show_category_name = True
    #     series.data_labels.font.size = Pt(8)
    #     series.data_labels.font.bold = True
    #     series.data_labels.font.name = 'Arial'
    #     # series.data_labels.show_series_name = False
    #     # series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    
    # # data table
    # _chart.has_data_table = True
    # # _chart.data_table.has_border_horizontal = True  # Add horizontal borders
    # # _chart.data_table.has_border_vertical = True  # Add vertical borders
    # # _chart.data_table.has_border_outline = True  # Add outline border

