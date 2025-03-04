import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import streamlit as st

def load_data(csv_file):
    return pd.read_csv(csv_file)

def get_slide_layout(prs, chart_layout_name):

    for layout in prs.slide_layouts:
        for placeholder in layout.placeholders:
            if placeholder.name == 'Chart Placeholder 8':
                return layout, placeholder.placeholder_format.idx
    return None


def get_data_from_csv():
    _data = st.file_uploader('Upload File', type=['csv'])
    return pd.read_csv(_data)

def load_ppt_template():
    selected_template = st.selectbox('TEMPLATES', options=['Template-1', 'Template-2'])

    if selected_template == 'Template-1':
        prs = Presentation('Templates/Template-1.pptx')
    
    if selected_template == 'Template-2':
        prs = Presentation('Templates/Template-2.pptx')

    return prs


if __name__ == '__main__':
    
    # load data
    df = get_data_from_csv()
    with st.expander('VIEW DATAFRAME'):
        st.dataframe(df)
    exit()
    # load template
    col11, col12 = st.columns([0.25, 0.75])
    with col11:
        prs = load_ppt_template()
    exit()
    col21, col22 = st.columns([0.25, 0.75])    
    # with col21:
        # btn_template = st.button('SELECT', use_container_width=True)

    _layouts = {}
    _layout_names = []
    # if btn_template:
    for layout in prs.slide_layouts:
        _layouts[layout.name] = layout
        _layout_names.append(layout.name)
        # st.write(layout.name)
            # for placeholder in layout.placeholders:
            #     st.write(f'{placeholder.name} - {placeholder.placeholder_format.idx}')
    
    # st.write(_layouts)
    # st.write(_layout_names)
    exit()

    with col12:
        selected_layout = st.selectbox('LAYOUTS', options=_layout_names)
        my_layout = _layouts[selected_layout]
        for placeholder in my_layout.placeholders:
            st.write(f'{placeholder.name} - {placeholder.placeholder_format.idx}')

    
    # slide_layout, _placeholderindex = get_slide_layout(prs, 'SOV')
    # slide = prs.slides.add_slide(slide_layout)
    
    # _title = slide.shapes.title
    # _title.text = 'My Sample Chart'

    # # Add Category
    # chart_data = CategoryChartData()
    # chart_data.categories = df['Date']

    # # Add Series
    # df_headers = df.columns.tolist()[1:]
    # for df_header in df_headers:
    #     chart_data.add_series(df_header, df[df_header])    
    
    # # Insert chart to the placeholder
    # chart_placeholder = slide.placeholders[_placeholderindex]
    # chart_frame = chart_placeholder.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data)
    # chart = chart_frame.chart

    # # Customize chart title
    # chart.has_title = True
    # chart.chart_title.text_frame.text = "Motorcycle Sales Over Time"
    # try:
    #     chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
    #     chart.chart_title.text_frame.paragraphs[0].font.name = 'Arial'
    # except:
    #     pass

    # # Customize category axis
    # chart.category_axis.axis_title.text_frame.text = "Date"
    # try:
    #     chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(8)
    #     chart.category_axis.axis_title.text_frame.paragraphs[0].font.name = 'Arial'
    # except:
    #     pass

    # # Customize category axis tick labels
    #     tick_labels_xml = chart.category_axis._element.find(".//c:txPr", namespaces={"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"})
        
    #     if tick_labels_xml is not None:
    #         print('Yes')


    #     chart.category_axis.tick_labels.font.bold = True
    #     chart.category_axis.tick_labels.font.color.rgb = RGBColor(100, 0, 0)
    #     chart.value_axis.tick_labels.rotation = 90

    # # Customize series axis
    # chart.value_axis.axis_title.text_frame.text = "Sales Count"
    # try:
    #     chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(8)
    #     chart.value_axis.axis_title.text_frame.paragraphs[0].font.name = 'Arial'
    # except:
    #     pass

    # # Customize gap between bars
    # chart.plots[0].gap_width = 50

    # # Apply custom colors to series
    # series_colors = [RGBColor(255, 0, 0), RGBColor(0, 128, 0), RGBColor(0, 0, 255), RGBColor(255, 165, 0), RGBColor(128, 0, 128)]

    # # chart transparency
    # # chart_frame.fill.solid()
    # # chart_frame.fill.fore_color.transparency = 0.0
    # # chart.plot_area.format.fill.solid()
    # # chart.plot_area.format.fill.fore_color.transparency = 0.0

    # # Customize legend
    # chart.has_legend = True
    # chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    # chart.legend.include_in_layout = False
    # chart.legend.font.size = Pt(8)
    # chart.legend.font.bold = True
    # chart.legend.font.name = 'Arial'
    
    

    # for layout in prs.slide_layouts:
    #     try:
    #         prs.slide_layouts.remove(layout)
    #     except:
    #         pass
    
    # prs.save('Output/output_presentation2.pptx')