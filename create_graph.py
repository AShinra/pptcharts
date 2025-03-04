import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt


def load_data(csv_file):
    return pd.read_csv(csv_file)

def load_ppt_template(ppt_file):
    return Presentation(ppt_file)

def get_slide_layout(prs, chart_layout_name):

    for layout in prs.slide_layouts:
        if layout.name == chart_layout_name:
            for placeholder in layout.placeholders:
                print(f'{placeholder.name}')
                if placeholder.name == 'Chart Placeholder':
                    print('found')
                    print(f'{placeholder.name}')
                    _placeholderindex = placeholder.placeholder_format.idx
                    print(f'{layout.name} --- {_placeholderindex}')
                    # return layout, _placeholderindex
    


    # return None




if __name__ == '__main__':

    df = load_data('Sample Data/Bar Data.csv')
    print(df)

    prs = load_ppt_template('Templates/chart_templates.pptx')

    # slide_layout, _placeholderindex = get_slide_layout(prs, 'Bar Chart Layout')
    get_slide_layout(prs, 'Bar Chart Layout')
    # slide = prs.slides.add_slide(slide_layout)

    # print(f'{slide_layout.name} -----')
    # print(_placeholderindex)
    exit()
    _title = slide.shapes.title
    _title.text = 'My Sample Chart'

    # Add Category
    chart_data = CategoryChartData()
    chart_data.categories = df['Date']

    # Add Series
    df_headers = df.columns.tolist()[1:]
    for df_header in df_headers:
        chart_data.add_series(df_header, df[df_header])    
    # chart_data.add_series('Honda', df['Honda'])
    # chart_data.add_series('Yamaha', df['Yamaha'])
    # chart_data.add_series('Suzuki', df['Suzuki'])
    # chart_data.add_series('Kawasaki', df['Kawasaki'])
    # chart_data.add_series('Motorstar', df['Motorstar'])

    # Insert chart to the placeholder
    chart_placeholder = slide.placeholders[_placeholderindex]
    chart_frame = chart_placeholder.insert_chart(XL_CHART_TYPE.COLUMN_STACKED, chart_data)
    chart = chart_frame.chart

    # Customize chart title and axis labels
    chart.has_title = True
    chart.chart_title.text_frame.text = "Motorcycle Sales Over Time"
    # chart.chart_title.text_frame.font.size = Pt(20)
    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = "Date"
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "Sales Count"


    # Customize legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.horz_offset = 0
    
    

    for layout in prs.slide_layouts:
        try:
            prs.slide_layouts.remove(layout)
        except:
            pass
    
    prs.save('Output/output_presentation.pptx')